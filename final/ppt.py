from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()

# Define slide layout and background color
slide_layout = prs.slide_layouts[5]  # Title Slide layout
background_color = RGBColor(10, 24, 74)  # Dark blue color

questions = [
    ("Which phenomenon describes the splitting of a single particle into two particles, such as electron-positron pairs, in the presence of a strong electric field?",
     ["A. Annihilation", "B. Photoelectric effect", "C. Compton scattering", "D. Pair production"]),
    ("Which was the first general-purpose electronic digital computer?",
     ["A. ENIAC", "B. UNIVAC", "C. IBM 701", "D. Mark I"]),
    ("Who is considered the father of computer science?",
     ["A. Alan Turing", "B. Charles Babbage", "C. John von Neumann", "D. Ada Lovelace"]),
    ("Which of the following is responsible for mediating the electromagnetic force?",
     ["A. Graviton", "B. Gluon", "C. Photon", "D. W boson"]),
    ("Which country is credited with the invention of paper?",
     ["A. India", "B. Greece", "C. China", "D. Egypt"]),
    ("What is the term for the hypothetical region of spacetime with extremely strong gravitational effects from which nothing, not even light, can escape?",
     ["A. White hole", "B. Event horizon", "C. Singularity", "D. Black hole"]),
    ("How many color dots make up one color pixel on a screen?",
     ["A. 265", "B. 16", "C. 8", "D. 3"]),
    ("What was the name of the mission that successfully landed the first humans on the Moon in 1969?",
     ["A. Apollo 10", "B. Apollo 13", "C. Apollo 8", "D. Apollo 11"]),
    ("Who is credited with the discovery of X-rays in 1895?",
     ["A. Max Planck", "B. Marie Curie", "C. Ernest Rutherford", "D. Wilhelm Conrad Röntgen"]),
    ("What is the term for the first successful test of the atomic bomb, conducted by the United States in July 1945?",
     ["A. Manhattan Project", "B. Operation Overlord", "C. Trinity test", "D. Hiroshima bombing"]),
    ("Who is the father of internet?",
     ["A. Charles Babbage", "B. Vint Cerf", "C. Dennis Ritchie", "D. Martin Cooper"]),
    ("When was the first public demonstration of electric light by Thomas Edison?",
     ["A. 1876", "B. 1879", "C. 1882", "D. 1885"]),
    ("What is the significance of the work of Leonardo da Vinci in mechanical engineering?",
     ["A. He invented the steam engine", "B. He designed numerous mechanical devices and concepts", "C. He developed the first hydraulic press", "D. He created the first working airplane"]),
    ("Program designed to perform specific task is known as",
     ["A. System Software", "B. Application Software", "C. Utility program", "D. Operating System"]),
    ("Bit stands for",
     ["A. Binary digits", "B. Bit of system", "C. A part of Byte", "D. All of the above"]),
]

# Add a slide for each question
for question, options in questions:
    slide = prs.slides.add_slide(slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(5.5))
    text_frame = textbox.text_frame
    
    # Set question text
    p = text_frame.add_paragraph()
    p.text = question
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White color
    
    # Set options text
    for option in options:
        p = text_frame.add_paragraph()
        p.text = option
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Save the presentation
ppt_file = "questions_presentation.pptx"
prs.save(ppt_file)

ppt_file
